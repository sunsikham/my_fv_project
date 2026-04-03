#!/usr/bin/env python3
"""Generate Q1 inside-A feature-mechanism figures for a 2-slide deck."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LinearSegmentedColormap, TwoSlopeNorm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path("/home/sunsik/my_fv_project")
Q1_DIR = PROJECT_ROOT / "results_fv" / "relation_condition_qwise" / "relA_relationA_ex__relB_relationB_ex__hd011861bcc" / "Q1"
OUT_DIR = PROJECT_ROOT / "docs" / "multi_feature_reweighting" / "slides" / "q1_inside_a_feature_mechanism"

REF_NAME = "AAA_ref"
FEATURE_NAMES = ["g0", "g1", "g2", "g3", "g4"]
STEP_LABELS = ["A1", "A2", "A3", "A4", "Aq"]
FEATURE_COLORS = {
    "g0": "#D9893D",
    "g1": "#7A6E9C",
    "g2": "#2F6BDA",
    "g3": "#0F9D8B",
    "g4": "#7EB6FF",
}


def _load_stepwise_arrays() -> tuple[dict[str, object], np.lib.npyio.NpzFile]:
    meta = json.loads((Q1_DIR / "_stepwise_a_states" / "stepwise_a_states_meta.json").read_text())
    npz = np.load(Q1_DIR / "_stepwise_a_states" / f"stepwise_a_states_{REF_NAME}.npz")
    return meta, npz


def _normalize(vec: np.ndarray) -> np.ndarray:
    vec = vec.astype(np.float32, copy=False)
    norm = float(np.linalg.norm(vec))
    if norm == 0.0:
        return np.zeros_like(vec)
    return (vec / norm).astype(np.float32, copy=False)


def _project_coeffs(v: np.ndarray, mu: np.ndarray, G: np.ndarray) -> np.ndarray:
    return (G.T @ (v.astype(np.float32, copy=False) - mu)).astype(np.float32, copy=False)


def _contrib_pr(contrib: np.ndarray) -> float:
    abs_sum = float(np.sum(np.abs(contrib)))
    sq_sum = float(np.sum(np.square(contrib)))
    if sq_sum == 0.0:
        return 0.0
    return float((abs_sum**2) / sq_sum)


def _build_top5_matched_bundle() -> dict[str, np.ndarray]:
    meta, z = _load_stepwise_arrays()
    aaa = z["Q1__AAA__matched__sum"].astype(np.float32)
    bab = z["Q1__BABA__matched__sum"].astype(np.float32)
    dad = z["Q1__DADA__matched__sum"].astype(np.float32)

    qA = np.load(Q1_DIR / "_vectors" / f"trial_vectors_{REF_NAME}_AAA.npy").astype(np.float32)
    qB = np.load(Q1_DIR / "_vectors" / f"trial_vectors_{REF_NAME}_BBB.npy").astype(np.float32)
    qD = np.load(Q1_DIR / "_vectors" / f"trial_vectors_{REF_NAME}_DDD.npy").astype(np.float32)
    a_anchor = qA.mean(axis=0).astype(np.float32, copy=False)
    b_anchor = qB.mean(axis=0).astype(np.float32, copy=False)
    d_anchor = qD.mean(axis=0).astype(np.float32, copy=False)
    b_minus_a = (b_anchor - a_anchor).astype(np.float32, copy=False)
    d_minus_a = (d_anchor - a_anchor).astype(np.float32, copy=False)

    rows = aaa.reshape(aaa.shape[0] * aaa.shape[1], aaa.shape[2]).astype(np.float32, copy=False)
    mu = rows.mean(axis=0, dtype=np.float64).astype(np.float32)
    centered = rows - mu
    _u, svals, vt = np.linalg.svd(centered, full_matrices=False)
    tol = max(centered.shape) * np.finfo(np.float32).eps * (svals[0] if len(svals) else 1.0)
    rank = int(np.sum(svals > tol))
    basis_rows = vt[:rank].astype(np.float32, copy=False)
    for idx in range(basis_rows.shape[0]):
        if float(np.dot(basis_rows[idx], b_minus_a)) < 0.0:
            basis_rows[idx] *= -1.0

    G = basis_rows[:5].T.astype(np.float32, copy=False)
    P = G @ G.T
    uB = _normalize(P @ b_minus_a)
    uD = _normalize(P @ d_minus_a)
    aB = (G.T @ uB).astype(np.float32, copy=False)
    aD = (G.T @ uD).astype(np.float32, copy=False)

    ids = meta["trial_ids_by_condition"]
    aaa_ids = ids["AAA"]
    bab_ids = ids["BABA"]
    dad_ids = ids["DADA"]
    aaa_map = {tid: i for i, tid in enumerate(aaa_ids)}
    bab_map = {tid: i for i, tid in enumerate(bab_ids)}
    dad_map = {tid: i for i, tid in enumerate(dad_ids)}
    common_bab = sorted(set(aaa_ids) & set(bab_ids))
    common_dad = sorted(set(aaa_ids) & set(dad_ids))

    mean_dc_bab = []
    mean_dc_dad = []
    mean_contrib_bab = []
    mean_contrib_dad = []
    T_B = []
    T_D = []
    pr_bab = []
    pr_dad = []

    for step_idx in range(len(STEP_LABELS)):
        dc_bab_list = []
        dc_dad_list = []
        contrib_bab_list = []
        contrib_dad_list = []
        for tid in common_bab:
            vA = aaa[aaa_map[tid], step_idx]
            vM = bab[bab_map[tid], step_idx]
            dc = _project_coeffs(vM, mu, G) - _project_coeffs(vA, mu, G)
            contrib = dc * aB
            dc_bab_list.append(dc)
            contrib_bab_list.append(contrib)
            T_B.append((step_idx, float(np.sum(contrib))))
            pr_bab.append((step_idx, _contrib_pr(contrib)))
        for tid in common_dad:
            vA = aaa[aaa_map[tid], step_idx]
            vM = dad[dad_map[tid], step_idx]
            dc = _project_coeffs(vM, mu, G) - _project_coeffs(vA, mu, G)
            contrib = dc * aD
            dc_dad_list.append(dc)
            contrib_dad_list.append(contrib)
            T_D.append((step_idx, float(np.sum(contrib))))
            pr_dad.append((step_idx, _contrib_pr(contrib)))

        mean_dc_bab.append(np.mean(np.stack(dc_bab_list), axis=0))
        mean_dc_dad.append(np.mean(np.stack(dc_dad_list), axis=0))
        mean_contrib_bab.append(np.mean(np.stack(contrib_bab_list), axis=0))
        mean_contrib_dad.append(np.mean(np.stack(contrib_dad_list), axis=0))

    T_B_mean = np.array(
        [np.mean([val for idx, val in T_B if idx == step_idx]) for step_idx in range(len(STEP_LABELS))],
        dtype=np.float32,
    )
    T_D_mean = np.array(
        [np.mean([val for idx, val in T_D if idx == step_idx]) for step_idx in range(len(STEP_LABELS))],
        dtype=np.float32,
    )
    PR_B_mean = np.array(
        [np.mean([val for idx, val in pr_bab if idx == step_idx]) for step_idx in range(len(STEP_LABELS))],
        dtype=np.float32,
    )
    PR_D_mean = np.array(
        [np.mean([val for idx, val in pr_dad if idx == step_idx]) for step_idx in range(len(STEP_LABELS))],
        dtype=np.float32,
    )

    return {
        "aB": aB,
        "aD": aD,
        "mean_dc_bab": np.stack(mean_dc_bab, axis=1),
        "mean_dc_dad": np.stack(mean_dc_dad, axis=1),
        "mean_contrib_bab": np.stack(mean_contrib_bab, axis=1),
        "mean_contrib_dad": np.stack(mean_contrib_dad, axis=1),
        "T_B_mean": T_B_mean,
        "T_D_mean": T_D_mean,
        "PR_B_mean": PR_B_mean,
        "PR_D_mean": PR_D_mean,
    }


def _draw_heatmap(ax, mat: np.ndarray, title: str, cmap, norm, show_ylabels: bool) -> None:
    im = ax.imshow(mat, aspect="auto", cmap=cmap, norm=norm)
    ax.set_title(title, fontsize=15, fontweight="bold")
    ax.set_xticks(np.arange(len(STEP_LABELS)))
    ax.set_xticklabels(STEP_LABELS, fontsize=11)
    ax.set_yticks(np.arange(len(FEATURE_NAMES)))
    ax.set_yticklabels(FEATURE_NAMES if show_ylabels else [""] * len(FEATURE_NAMES), fontsize=12)
    for i in range(mat.shape[0]):
        for j in range(mat.shape[1]):
            ax.text(
                j,
                i,
                f"{mat[i, j]:+.2f}",
                ha="center",
                va="center",
                fontsize=10,
                color="black",
                fontweight="bold" if abs(mat[i, j]) >= 0.45 else None,
            )
    ax.set_xticks(np.arange(-0.5, len(STEP_LABELS), 1), minor=True)
    ax.set_yticks(np.arange(-0.5, len(FEATURE_NAMES), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=1.5)
    ax.tick_params(which="minor", bottom=False, left=False)
    for spine in ax.spines.values():
        spine.set_visible(False)
    return im


def _draw_alignment_panel(ax, payload: dict[str, np.ndarray]) -> None:
    ax.axis("off")
    y0 = 0.88
    dy = 0.16
    for idx, _g in enumerate(FEATURE_NAMES):
        g_name = f"g{idx}"
        y = y0 - idx * dy
        ax.text(
            0.0,
            y,
            g_name,
            fontsize=12.5,
            fontweight="bold",
            color=FEATURE_COLORS[g_name],
            va="top",
            ha="left",
            transform=ax.transAxes,
        )
        ax.text(
            0.18,
            y,
            f"align_B = {payload['aB'][idx]:+.3f}",
            fontsize=10.5,
            family="DejaVu Sans Mono",
            color="#0f172a",
            va="top",
            ha="left",
            transform=ax.transAxes,
        )
        ax.text(
            0.18,
            y - 0.055,
            f"align_D = {payload['aD'][idx]:+.3f}",
            fontsize=10.5,
            family="DejaVu Sans Mono",
            color="#0f172a",
            va="top",
            ha="left",
            transform=ax.transAxes,
        )


def _plot_feature_change_figure(payload: dict[str, np.ndarray]) -> Path:
    fig = plt.figure(figsize=(16.2, 6.4))
    gs = fig.add_gridspec(1, 3, width_ratios=[0.9, 1.35, 1.35], wspace=0.04)
    ax_info = fig.add_subplot(gs[0, 0])
    ax_bab = fig.add_subplot(gs[0, 1])
    ax_dad = fig.add_subplot(gs[0, 2])
    cmap = LinearSegmentedColormap.from_list("delta_c", ["#2F6BDA", "#F7F7F7", "#D9534F"])
    all_vals = np.concatenate([payload["mean_dc_bab"].ravel(), payload["mean_dc_dad"].ravel()])
    vmax = float(np.max(np.abs(all_vals)))
    norm = TwoSlopeNorm(vmin=-vmax, vcenter=0.0, vmax=vmax)

    _draw_alignment_panel(ax_info, payload)
    _draw_heatmap(ax_bab, payload["mean_dc_bab"], "BABABA", cmap, norm, True)
    im = _draw_heatmap(ax_dad, payload["mean_dc_dad"], "DADADA", cmap, norm, False)

    cax = fig.add_axes([0.93, 0.18, 0.016, 0.62])
    cbar = fig.colorbar(im, cax=cax)
    cbar.set_label("mean Δc", fontsize=11)

    fig.subplots_adjust(left=0.05, right=0.91, top=0.90, bottom=0.12)

    out_path = OUT_DIR / "q1_inside_a_feature_changes.png"
    fig.savefig(out_path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return out_path


def _stack_signed_bars(ax, contrib: np.ndarray, totals: np.ndarray, color_order: list[str]) -> None:
    x = np.arange(contrib.shape[1])
    pos_bottom = np.zeros_like(x, dtype=np.float32)
    neg_bottom = np.zeros_like(x, dtype=np.float32)
    for idx, g_name in enumerate(color_order):
        vals = contrib[idx]
        pos = np.clip(vals, 0.0, None)
        neg = np.clip(vals, None, 0.0)
        ax.bar(x, pos, bottom=pos_bottom, color=FEATURE_COLORS[g_name], width=0.68, label=g_name)
        ax.bar(x, neg, bottom=neg_bottom, color=FEATURE_COLORS[g_name], width=0.68)
        pos_bottom += pos
        neg_bottom += neg

    ax.plot(x, totals, color="black", marker="o", linewidth=2.2, label="total inside-A push")
    ax.axhline(0.0, color="#555555", linewidth=1.0)
    ax.set_xticks(x)
    ax.set_xticklabels(STEP_LABELS, fontsize=11)
    ax.set_ylabel("signed contribution", fontsize=11)
    ax.set_xlabel("Matched A slots", fontsize=11)
    ax.grid(axis="y", alpha=0.25)


def _plot_contribution_figure(
    contrib: np.ndarray,
    totals: np.ndarray,
    *,
    out_name: str,
) -> Path:
    fig, ax = plt.subplots(1, 1, figsize=(12.6, 5.8))
    color_order = ["g0", "g1", "g2", "g3", "g4"]

    _stack_signed_bars(
        ax,
        contrib,
        totals,
        color_order,
    )
    fig.subplots_adjust(left=0.10, right=0.98, top=0.97, bottom=0.16)

    out_path = OUT_DIR / out_name
    fig.savefig(out_path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return out_path


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    payload = _build_top5_matched_bundle()
    feature_change_path = _plot_feature_change_figure(payload)
    contribution_bab_path = _plot_contribution_figure(
        payload["mean_contrib_bab"],
        payload["T_B_mean"],
        out_name="q1_inside_a_feature_contributions_bab.png",
    )
    contribution_dad_path = _plot_contribution_figure(
        payload["mean_contrib_dad"],
        payload["T_D_mean"],
        out_name="q1_inside_a_feature_contributions_dad.png",
    )
    pptx_path = _build_pptx(feature_change_path, contribution_bab_path, contribution_dad_path, payload)
    metrics = {
        "feature_alignments": {
            f"g{idx}": {
                "align_B": float(payload["aB"][idx]),
                "align_D": float(payload["aD"][idx]),
            }
            for idx in range(5)
        },
        "T_B_mean": [float(x) for x in payload["T_B_mean"]],
        "T_D_mean": [float(x) for x in payload["T_D_mean"]],
        "PR_B_mean": [float(x) for x in payload["PR_B_mean"]],
        "PR_D_mean": [float(x) for x in payload["PR_D_mean"]],
        "feature_change_png": str(feature_change_path),
        "feature_contribution_bab_png": str(contribution_bab_path),
        "feature_contribution_dad_png": str(contribution_dad_path),
        "pptx": str(pptx_path),
    }
    (OUT_DIR / "q1_inside_a_feature_metrics.json").write_text(
        json.dumps(metrics, ensure_ascii=True, indent=2),
        encoding="utf-8",
    )


def _add_textbox(slide, left, top, width, height, text, *, font_size=20, bold=False, color=(15, 23, 42), align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return tx


def _add_takeaway_box(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(6.65),
        Inches(12.2),
        Inches(0.48),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(247, 250, 252)
    shape.line.color.rgb = RGBColor(203, 213, 225)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = text
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.bold = False
    run.font.color.rgb = RGBColor(15, 23, 42)


def _add_color_legend(slide, *, left: float, top: float, include_total: bool = False) -> None:
    items = [("g0", FEATURE_COLORS["g0"]), ("g1", FEATURE_COLORS["g1"]), ("g2", FEATURE_COLORS["g2"]), ("g3", FEATURE_COLORS["g3"]), ("g4", FEATURE_COLORS["g4"])]
    if include_total:
        items.append(("total", "#000000"))
    x = left
    for label, color in items:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(top),
            Inches(0.16),
            Inches(0.12),
        )
        box.fill.solid()
        rgb = tuple(int(color[i : i + 2], 16) for i in (1, 3, 5))
        box.fill.fore_color.rgb = RGBColor(*rgb)
        box.line.fill.background()
        _add_textbox(
            slide,
            Inches(x + 0.20),
            Inches(top - 0.03),
            Inches(0.38),
            Inches(0.16),
            label,
            font_size=11,
            bold=True,
            color=rgb,
        )
        x += 0.72


def _add_vertical_legend(slide, *, left: float, top: float, include_total: bool = False) -> None:
    items = [("g0", FEATURE_COLORS["g0"]), ("g1", FEATURE_COLORS["g1"]), ("g2", FEATURE_COLORS["g2"]), ("g3", FEATURE_COLORS["g3"]), ("g4", FEATURE_COLORS["g4"])]
    if include_total:
        items.append(("total", "#000000"))
    y = top
    for label, color in items:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(y),
            Inches(0.16),
            Inches(0.12),
        )
        box.fill.solid()
        rgb = tuple(int(color[i : i + 2], 16) for i in (1, 3, 5))
        box.fill.fore_color.rgb = RGBColor(*rgb)
        box.line.fill.background()
        _add_textbox(
            slide,
            Inches(left + 0.22),
            Inches(y - 0.03),
            Inches(0.70),
            Inches(0.16),
            label,
            font_size=13,
            color=(31, 41, 55),
        )
        y += 0.28


def _build_pptx(feature_change_path: Path, contribution_bab_path: Path, contribution_dad_path: Path, payload: dict[str, np.ndarray]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.22),
        Inches(8.9),
        Inches(0.45),
        "Which Inside-A Features Change Over Steps?",
        font_size=25,
        bold=True,
    )
    _add_textbox(
        slide,
        Inches(0.57),
        Inches(0.67),
        Inches(11.8),
        Inches(0.38),
        "Inside-A carries most of the stepwise change. The next question is which A-features are being reweighted across matched A slots.",
        font_size=14,
        color=(51, 65, 85),
    )
    slide.shapes.add_picture(str(feature_change_path), Inches(0.45), Inches(1.0), width=Inches(12.4))
    _add_takeaway_box(
        slide,
        "Takeaway. Inside A, the same role-defined features do not stay fixed: g2/g3/g4 are amplified along the B branch, while the D branch strengthens g0 and suppresses anti-D axes such as g1 and g3.",
    )

    # Slide 2
    slide = prs.slides.add_slide(blank)
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.22),
        Inches(8.9),
        Inches(0.45),
        "BABABA: Which Features Drive the B-Directed Drift?",
        font_size=25,
        bold=True,
    )
    _add_textbox(
        slide,
        Inches(0.57),
        Inches(0.67),
        Inches(12.0),
        Inches(0.56),
        "Contribution asks not just which features changed, but which of those changes actually pushed the state toward B. Each colored segment shows how much one inside-A feature contributes to that B-directed drift, relative to the AAA baseline. Segments above zero support B-directed drift; segments below zero oppose it.",
        font_size=13,
        color=(51, 65, 85),
    )
    _add_color_legend(slide, left=0.62, top=1.18, include_total=True)
    slide.shapes.add_picture(str(contribution_bab_path), Inches(0.70), Inches(1.56), width=Inches(11.95))

    # Slide 3
    slide = prs.slides.add_slide(blank)
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.22),
        Inches(8.9),
        Inches(0.45),
        "DADADA: Which Features Drive the D-Directed Drift?",
        font_size=25,
        bold=True,
    )
    _add_textbox(
        slide,
        Inches(0.57),
        Inches(0.67),
        Inches(12.0),
        Inches(0.56),
        "Contribution asks not just which features changed, but which of those changes actually pushed the state toward D. Each colored segment shows how much one inside-A feature contributes to that D-directed drift, relative to the AAA baseline. Segments above zero support D-directed drift; segments below zero oppose it.",
        font_size=13,
        color=(51, 65, 85),
    )
    _add_color_legend(slide, left=0.62, top=1.18, include_total=True)
    slide.shapes.add_picture(str(contribution_dad_path), Inches(0.70), Inches(1.56), width=Inches(11.95))

    out_path = OUT_DIR / "q1_inside_a_feature_mechanism.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    main()
